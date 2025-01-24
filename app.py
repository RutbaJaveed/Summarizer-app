import streamlit as st  
import pymupdf   #For pdf
from docx import Document  # For docx
from pptx import Presentation  # For pptx
from dotenv import load_dotenv
import os
import google.generativeai as genai


load_dotenv()

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")


genai.configure(api_key=GOOGLE_API_KEY)


#Initialize the generative model
model = genai.GenerativeModel('gemini-1.5-flash')

#Function to interact with the model
def model_answer(ques, max_words):
    try:
        config = genai.GenerationConfig(max_output_tokens = max_words, temperature = 0.5)
        response = model.generate_content(ques, generation_config=config)
        summary = response.text

        #Ensure the summary ends with a complete sentence
        if summary and not summary.endswith('.'):
            last_period = summary.rfind('.')
            if last_period != -1:
                summary = summary[:last_period + 1]                
        return summary
    except Exception as e:
        return f"An error occurred: {e}"
    
#Function to extract text from PDF
def extract_pdf_content(uploaded_file):
    try:
        pdf_doc = pymupdf.open(stream=uploaded_file.read(),filetype="pdf")

        content_stream = ""

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            content_stream += page.get_text()

        pdf_doc.close()
        return content_stream
    
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return None

#Function to extract text from DOCX  
def extract_docx_content(uploaded_file):
    try:
        doc = Document(uploaded_file)
        content_stream = "\n".join([p.text for p in doc.paragraphs])
        return content_stream
    except Exception as e:
        st.error(f"An error occurred while processing the DOCX: {e}")
        return None
    

 #Function to extract text from PPTX   
def extract_pptx_content(uploaded_file):
    try:
        presentation = Presentation(uploaded_file)
        content_stream = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    content_stream += shape.text + "\n"
        return content_stream
    except Exception as e:
        st.error(f"An error occurred while processing the PPTX: {e}")
        return None

#Function to extract text from TXT   
def extract_txt_content(uploaded_file):
    try:
        content_stream = uploaded_file.read().decode("utf-8")
        return content_stream
    except Exception as e:
        st.error(f"An error occurred while processing the TXT: {e}")
        return None


#Streamlit App
st.set_page_config(page_title="Multi-doc Summarizer", page_icon="📄",layout="wide")
st.markdown(
    """
    <div style="background-color:#4A90E2;padding:5px;border-radius:5px">
    <h1 style="color:white;text-align:center;">📄 Multi-doc Summarizer</h1>
    </div><br>
    """,True
)

uploaded_file = st.file_uploader("__Upload a document__", type=["pdf", "doc", "docx", "ppt", "pptx", "txt"],label_visibility="collapsed")
summary_word_limit = st.slider("Choose Maximum Summary Word Count: ", min_value=100, max_value=2000)


if uploaded_file is not None:
    file_type = uploaded_file.type
    content = None

    # Handle file type
    if file_type == "application/pdf":
        content = extract_pdf_content(uploaded_file)
    elif file_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
        content = extract_docx_content(uploaded_file)
    elif file_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
        content = extract_pptx_content(uploaded_file)
    elif file_type == "text/plain":
        content = extract_txt_content(uploaded_file)
    else:
        st.error("Unsupported file format. Please upload a PDF, DOCX, PPTX, or TXT file.")

#If content is successfully extracted
    if content and st.button("Summarize Content"):
        base_prompt = (f"Summarize the following document in no more than {summary_word_limit} words." 
        "Ensure the summary is concise, captures the main ideas clearly, and ends with a complete sentence. "
        "If the word limit is insufficient to cover everything, prioritize the most critical points end ensure the summary reads coherently.")
        final_prompt = base_prompt + content
        summarized_content = model_answer(final_prompt,summary_word_limit) 
        if "An error occurred" in summarized_content:
                st.error("An issue occurred while generating the summary. Please try again.")
        else:
            st.markdown("### **Summarized Content**")
            st.write(summarized_content) 
    elif content is None and not st.button("Summarize Content"):
        st.warning("The document appears to be empty or could not be processed. Please try with another file.")